#!/usr/bin/env python3
"""
Gridiron Partners — Presentation Generator
Modifies the PPTX template directly with live Supabase data.
Static slides stay pixel-perfect; dynamic slides get updated data.

Usage:
  python generate_presentation.py
  python generate_presentation.py --client "Beaver County" --month 2025-08
  python generate_presentation.py --email you@email.com --password pw --client "BAC" --month 2025-12

Requirements:
  pip install python-pptx Pillow matplotlib
"""

import os, sys, json, math, argparse, calendar, getpass
from datetime import datetime
from collections import defaultdict
from io import BytesIO
import urllib.request, urllib.parse

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
import lxml.etree as etree

# ═══════════════════════════════════════════════════════════════════
# CONFIG
# ═══════════════════════════════════════════════════════════════════
SB_URL  = 'https://nzinvxticgyjobkqxxhl.supabase.co'
SB_KEY  = ('eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIs'
           'InJlZiI6Im56aW52eHRpY2d5am9ia3F4eGhsIiwicm9sZSI6ImFub24iLCJp'
           'YXQiOjE3Njg5NTUwOTIsImV4cCI6MjA4NDUzMTA5Mn0.YhaW1zesBaEC6lRM'
           'Wx6_gyZN5Uhh8FX2wgahunHQfok')

MONTH_NAMES = ['','January','February','March','April','May','June',
               'July','August','September','October','November','December']

FI_SECTORS = ['Diversified','Corp/HY','Preferreds','Floating Rate',
              'Municipals','MBB','Low Duration','Convertibles','Foreign']

SECTOR_MAP = {
    'Diversified':'Diversified',
    'Corp HY':'Corp/HY','Corp/HY':'Corp/HY','Corporate HY':'Corp/HY',
    'Corp High Yield':'Corp/HY','High Yield':'Corp/HY',
    'Preferred':'Preferreds','Preferreds':'Preferreds',
    'Floating Rate':'Floating Rate','Senior Loans':'Floating Rate',
    'Municipal':'Municipals','Municipals':'Municipals','Muni':'Municipals',
    'MBS':'MBB','MBB':'MBB','Mortgage':'MBB','Mortgage Backed':'MBB',
    'Low Duration':'Low Duration','Limited Duration':'Low Duration','Short Duration':'Low Duration',
    'Convertible':'Convertibles','Convertibles':'Convertibles','Convert':'Convertibles',
    'Foreign':'Foreign','Global':'Foreign','Emerging':'Foreign','EM':'Foreign',
    'Foreign/EM':'Foreign','Intl':'Foreign',
}

# Chart pair order must match slide order (slides 7-11)
CHART_PAIRS = [
    ('Diversified', 'Corp/HY'),
    ('Preferreds', 'Floating Rate'),
    ('Municipals', 'MBB'),
    ('Low Duration', 'Convertibles'),
    ('Foreign',),
]

# ═══════════════════════════════════════════════════════════════════
# SUPABASE HELPERS
# ═══════════════════════════════════════════════════════════════════
def sb_get(path, token=None):
    url = f"{SB_URL}/rest/v1/{path}"
    req = urllib.request.Request(url, headers={
        'apikey': SB_KEY,
        'Authorization': f'Bearer {token or SB_KEY}',
    })
    with urllib.request.urlopen(req) as r:
        return json.loads(r.read())

def sb_all(table, select='*', filt='', order='', token=None):
    rows, off = [], 0
    while True:
        p = f"{table}?select={select}&limit=10000&offset={off}"
        if filt: p += f"&{filt}"
        if order: p += f"&order={order}"
        b = sb_get(p, token)
        rows.extend(b)
        if len(b) < 10000: break
        off += 10000
    return rows

def sb_login(email, password):
    url = f"{SB_URL}/auth/v1/token?grant_type=password"
    data = json.dumps({'email': email, 'password': password}).encode()
    req = urllib.request.Request(url, data=data, headers={
        'apikey': SB_KEY, 'Content-Type': 'application/json',
    })
    try:
        with urllib.request.urlopen(req) as r:
            return json.loads(r.read())['access_token']
    except urllib.error.HTTPError as e:
        body = e.read().decode()
        print(f"\n  LOGIN FAILED ({e.code}): {body}")
        print("  Check your email/password and try again.")
        print("  Tip: use --skip-auth to try with the anon key instead.")
        sys.exit(1)

# ═══════════════════════════════════════════════════════════════════
# FORMATTING
# ═══════════════════════════════════════════════════════════════════
def fmt_mo(ym):
    y, m = ym.split('-'); return f"{MONTH_NAMES[int(m)]} {y}"

def fmt_mo_short(ym):
    y, m = ym.split('-'); return f"{MONTH_NAMES[int(m)][:3]} {y}"

def fmt_date(ym):
    y, m = int(ym[:4]), int(ym[5:])
    return f"{MONTH_NAMES[m]} {calendar.monthrange(y, m)[1]}, {y}"

def fmt_cat(ym):
    y, m = ym.split('-'); return f"{MONTH_NAMES[int(m)][:3]}-{y[2:]}"

def fmt_dollar(v):
    if v is None: return '$0.00'
    return f"${abs(v):,.2f}" if v >= 0 else f"-${abs(v):,.2f}"

def fmt_pct(v, dp=2):
    if v is None: return '—'
    return f"{v:.{dp}f}%"

# ═══════════════════════════════════════════════════════════════════
# DATA LOADING
# ═══════════════════════════════════════════════════════════════════
def load_all(token):
    print("  Loading clients..."); clients = sb_all('clients', token=token, order='client_name.asc')
    print("  Loading strategies..."); strats = sb_all('strategies', token=token)
    print("  Loading index returns..."); idx = sb_all('index_returns', token=token, order='month_end_date.asc')
    ir = defaultdict(dict)
    for r in idx: ir[r['index_id']][r['month_end_date'][:7]] = r['monthly_return']

    print("  Loading CEF tickers..."); tickers = sb_all('cef_tickers', select='ticker,sector,asset_class,fund_name', token=token)
    print("  Loading CEF daily (may take a moment)...")
    cef = sb_all('cef_daily', select='ticker,trade_date,discount_pct',
                 filt='discount_pct=not.is.null', order='trade_date.asc', token=token)
    print(f"    → {len(cef):,} rows")
    return dict(clients=clients, strategies=strats, index_returns=dict(ir),
                tickers=tickers, cef_daily=cef)

def load_monthly(cid, token):
    return sb_all('monthly_data', filt=f'client_id=eq.{cid}', order='month_end_date.asc', token=token)

# ═══════════════════════════════════════════════════════════════════
# SECTOR DISCOUNT DATA
# ═══════════════════════════════════════════════════════════════════
def build_sectors(tickers, cef_daily):
    tk_sec = {}
    for t in tickers:
        m = SECTOR_MAP.get(t['sector'], t['sector'])
        if m in FI_SECTORS: tk_sec[t['ticker']] = m

    tk_mo = defaultdict(dict)
    for d in cef_daily:
        s = tk_sec.get(d['ticker'])
        if s: tk_mo[d['ticker']][d['trade_date'][:7]] = d['discount_pct']

    months = sorted(set(ym for tk in tk_mo.values() for ym in tk))
    out = {}
    for sec in FI_SECTORS:
        stks = [t for t, s in tk_sec.items() if s == sec]
        series, vals = [], []
        for ym in months:
            v = [tk_mo[t][ym] for t in stks if ym in tk_mo[t]]
            if v:
                a = sum(v)/len(v); series.append((ym, a)); vals.append(a)
            else:
                series.append((ym, None))
        ha = sum(vals)/len(vals) if vals else 0
        cur = vals[-1] if vals else None
        out[sec] = dict(series=series, hist_avg=ha, current=cur)
    return out

# ═══════════════════════════════════════════════════════════════════
# PPTX MODIFICATION: SLIDE 1 (COVER)
# ═══════════════════════════════════════════════════════════════════
def update_cover(slide, client_name, month):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == 'TextBox 2':
            for p in shape.text_frame.paragraphs:
                txt = p.text.strip()
                if not txt or txt == 'Prepared For:': continue
                # Detect date line vs client name
                if any(mn in txt for mn in MONTH_NAMES[1:]) or (len(txt) < 30 and any(c.isdigit() for c in txt)):
                    for r in p.runs: r.text = fmt_mo(month)
                elif txt != 'Prepared For:':
                    for r in p.runs: r.text = client_name

# ═══════════════════════════════════════════════════════════════════
# PPTX MODIFICATION: SLIDES 7-11 (SECTOR CHARTS)
# ═══════════════════════════════════════════════════════════════════
def _update_chart_xml(chart, cats, s0_vals, s1_vals):
    """Replace area chart cache data."""
    cs = chart._chartSpace
    ac = cs.find(qn('c:chart')).find(qn('c:plotArea')).find(qn('c:areaChart'))
    if ac is None: return
    for s_idx, ser in enumerate(ac.findall(qn('c:ser'))):
        vals = s0_vals if s_idx == 0 else s1_vals
        # Values
        nc = ser.find(qn('c:val')).find(qn('c:numRef')).find(qn('c:numCache'))
        for pt in nc.findall(qn('c:pt')): nc.remove(pt)
        nc.find(qn('c:ptCount')).set('val', str(len(vals)))
        for i, v in enumerate(vals):
            pt = etree.SubElement(nc, qn('c:pt')); pt.set('idx', str(i))
            ve = etree.SubElement(pt, qn('c:v')); ve.text = str(v if v is not None else 0)
        # Categories
        sc = ser.find(qn('c:cat')).find(qn('c:strRef')).find(qn('c:strCache'))
        for pt in sc.findall(qn('c:pt')): sc.remove(pt)
        sc.find(qn('c:ptCount')).set('val', str(len(cats)))
        for i, c in enumerate(cats):
            pt = etree.SubElement(sc, qn('c:pt')); pt.set('idx', str(i))
            ve = etree.SubElement(pt, qn('c:v')); ve.text = c

def _update_stats_tbl(tbl, name, ha, cur):
    """Update 5×1 stats table."""
    for p in tbl.cell(0,0).text_frame.paragraphs:
        for r in p.runs: r.text = name
    for p in tbl.cell(2,0).text_frame.paragraphs:
        for r in p.runs: r.text = fmt_pct(ha*100)
    for p in tbl.cell(4,0).text_frame.paragraphs:
        for r in p.runs: r.text = fmt_pct(cur*100) if cur is not None else '—'

def update_sector_slides(prs, sec_data, month):
    for pi, pair in enumerate(CHART_PAIRS):
        slide = prs.slides[6 + pi]
        charts = sorted([s for s in slide.shapes if s.has_chart], key=lambda s: s.top)
        tables = sorted([s for s in slide.shapes if s.has_table], key=lambda s: s.top)

        for si, sec in enumerate(pair):
            sd = sec_data.get(sec)
            if not sd or si >= len(charts): continue
            cats = [fmt_cat(ym) for ym, _ in sd['series']]
            s0 = [sd['hist_avg']] * len(cats)
            s1 = [v if v is not None else 0 for _, v in sd['series']]
            _update_chart_xml(charts[si].chart, cats, s0, s1)
            if si < len(tables):
                _update_stats_tbl(tables[si].table, sec, sd['hist_avg'], sd['current'])
            print(f"    ✓ {sec}")

        # Date on slide 7
        if pi == 0:
            for s in slide.shapes:
                if s.has_text_frame and s.name == 'TextBox 13':
                    for p in s.text_frame.paragraphs:
                        for r in p.runs: r.text = fmt_mo(month)

# ═══════════════════════════════════════════════════════════════════
# PPTX MODIFICATION: SLIDE 12 (ASSET ALLOCATION)
# ═══════════════════════════════════════════════════════════════════
def update_allocation(slide, month):
    """Update date header in allocation table. Data stays as composite default
       until per-client allocation data is available."""
    for s in slide.shapes:
        if s.has_table:
            tbl = s.table
            # Update first date column header
            for p in tbl.cell(1, 2).text_frame.paragraphs:
                for r in p.runs: r.text = fmt_mo_short(month)
            break

# ═══════════════════════════════════════════════════════════════════
# PPTX MODIFICATION: SLIDE 13 (PERFORMANCE STATEMENT IMAGE)
# ═══════════════════════════════════════════════════════════════════
def _perf_calcs(data, month):
    """Calculate all performance numbers."""
    if not data:
        print("    WARNING: No monthly data")
        return None
    year = month[:4]
    md = next((d for d in data if d['month_end_date'][:7] == month), None)
    if not md: return None

    mo = int(month[5:])
    qs = f"{year}-{'01' if mo<=3 else '04' if mo<=6 else '07' if mo<=9 else '10'}"
    ytd_d = [d for d in data if d['month_end_date'][:4]==year and d['month_end_date'][:7]<=month]
    qtd_d = [d for d in data if qs<=d['month_end_date'][:7]<=month]
    si_d  = [d for d in data if d['month_end_date'][:7]<=month]

    beg_q = qtd_d[0]['beginning_balance'] if qtd_d else md['beginning_balance']
    beg_y = ytd_d[0]['beginning_balance'] if ytd_d else md['beginning_balance']
    add_q = sum(d.get('contributions') or 0 for d in qtd_d)
    add_y = sum(d.get('contributions') or 0 for d in ytd_d)
    wd_q  = sum(d.get('withdrawals') or 0 for d in qtd_d)
    wd_y  = sum(d.get('withdrawals') or 0 for d in ytd_d)
    chg_m = (md['ending_balance'] or 0) - (md['beginning_balance'] or 0) - (md.get('contributions') or 0) + (md.get('withdrawals') or 0)
    chg_q = (md['ending_balance'] or 0) - (beg_q or 0) - add_q + wd_q
    chg_y = (md['ending_balance'] or 0) - (beg_y or 0) - add_y + wd_y

    # Returns stored as decimals (0.011 = 1.1%)
    gm = (md.get('gross_return_pct') or 0) * 100
    nm = (md.get('net_return_pct') or 0) * 100
    fm = gm - nm

    def chain(lst, k):
        c=1
        for d in lst: c *= (1 + (d.get(k) or 0))
        return c - 1

    ytd_g = chain(ytd_d,'gross_return_pct'); ytd_n = chain(ytd_d,'net_return_pct')
    qtd_g = chain(qtd_d,'gross_return_pct'); qtd_n = chain(qtd_d,'net_return_pct')
    si_g  = chain(si_d,'gross_return_pct');  si_n  = chain(si_d,'net_return_pct')
    n = len(si_d)
    ann_n = (math.pow(1+si_n, 12/n)-1) if n>12 else si_n
    fm1 = data[0]['month_end_date'][:7] if data else month

    return dict(md=md, beg_q=beg_q, beg_y=beg_y, add_q=add_q, add_y=add_y,
                wd_q=wd_q, wd_y=wd_y, chg_m=chg_m, chg_q=chg_q, chg_y=chg_y,
                gm=gm, nm=nm, fm=fm, ytd_g=ytd_g, ytd_n=ytd_n,
                qtd_g=qtd_g, qtd_n=qtd_n, si_g=si_g, si_n=si_n,
                ann_n=ann_n, first_mo=fm1, n_months=n)


def gen_perf_image(client_name, strat_name, month, data):
    """Render performance statement as a high-res PNG."""
    import matplotlib; matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    from matplotlib.patches import Rectangle

    c = _perf_calcs(data, month)
    if not c: return None
    md = c['md']
    GOLD, DK, GRN, RED = '#CAA879','#1A1A2E','#059669','#DC2626'
    def pc(v): return GRN if v>=0 else RED

    fig, ax = plt.subplots(figsize=(8.54, 5.2))
    ax.set_xlim(0, 8.54); ax.set_ylim(0, 5.2); ax.axis('off')
    fig.patch.set_facecolor('white')

    y = 5.0
    ax.text(0.15, y, client_name, fontsize=13, fontweight='bold', color=DK, va='top', fontfamily='sans-serif')
    y -= 0.24
    ax.text(0.15, y, strat_name, fontsize=8.5, fontweight='bold', color=GOLD, va='top')
    y -= 0.17
    ax.text(0.15, y, fmt_date(month), fontsize=7.5, color='#888', va='top')
    y -= 0.32

    cx = [0.15, 3.6, 5.5, 7.35]  # column left edges
    rw = [3.35, 1.8, 1.8, 1.19]  # column widths

    def draw_header(labels, yp):
        for i, (lbl, w) in enumerate(zip(labels, rw)):
            ax.add_patch(Rectangle((cx[i], yp-0.01), w, 0.22, fc=GOLD, ec='none'))
            ha = 'left' if i==0 else 'right'
            xp = cx[i]+0.05 if i==0 else cx[i]+w-0.05
            ax.text(xp, yp+0.10, lbl, fontsize=6.5, fontweight='bold', color='white', va='center', ha=ha)
        return yp - 0.26

    def draw_row(lbl, vals, yp, bold=False):
        fw = 'bold' if bold else 'normal'
        if bold:
            ax.plot([cx[0], cx[3]+rw[3]], [yp+0.12, yp+0.12], color='#aaa', lw=0.6)
        ax.text(cx[0]+0.05, yp, lbl, fontsize=7.5, color='#333', va='center', fontweight=fw)
        for i, v in enumerate(vals):
            xp = cx[i+1]+rw[i+1]-0.05
            if isinstance(v, tuple):  # (text, color)
                ax.text(xp, yp, v[0], fontsize=7.5, color=v[1], va='center', ha='right', fontweight=fw)
            else:
                ax.text(xp, yp, v, fontsize=7.5, color='#333', va='center', ha='right', fontweight=fw)
        ax.plot([cx[0], cx[3]+rw[3]], [yp-0.10, yp-0.10], color='#e8e8e8', lw=0.4)
        return yp - 0.22

    # Account Summary
    y = draw_header(['ACCOUNT SUMMARY','Month-to-Date','Quarter-to-Date','Year-to-Date'], y)
    y = draw_row('Beginning Balance', [fmt_dollar(md['beginning_balance']), fmt_dollar(c['beg_q']), fmt_dollar(c['beg_y'])], y)
    y = draw_row('Additions', [fmt_dollar(md.get('contributions') or 0), fmt_dollar(c['add_q']), fmt_dollar(c['add_y'])], y)
    y = draw_row('Withdrawals', [fmt_dollar(md.get('withdrawals') or 0), fmt_dollar(c['wd_q']), fmt_dollar(c['wd_y'])], y)
    y = draw_row('Change in Value', [fmt_dollar(c['chg_m']), fmt_dollar(c['chg_q']), fmt_dollar(c['chg_y'])], y)
    y = draw_row('Ending Value', [fmt_dollar(md['ending_balance'])]*3, y, bold=True)
    y -= 0.18

    # Performance Summary
    y = draw_header(['PERFORMANCE SUMMARY','Month-to-Date','Quarter-to-Date','Year-to-Date'], y)
    y = draw_row('Gross Performance',
                 [(fmt_pct(c['gm']), pc(c['gm'])), (fmt_pct(c['qtd_g']*100), pc(c['qtd_g'])),
                  (fmt_pct(c['ytd_g']*100), pc(c['ytd_g']))], y)
    fee_m_d = c['fm'] * (md['beginning_balance'] or 0) / 100 if c['fm']>0 else 0
    fee_q_d = (c['qtd_g']-c['qtd_n'])*(c['beg_q'] or 0)
    fee_y_d = (c['ytd_g']-c['ytd_n'])*(c['beg_y'] or 0)
    y = draw_row('Management Fee', [fmt_dollar(fee_m_d), fmt_dollar(fee_q_d), fmt_dollar(fee_y_d)], y)
    y = draw_row('Net Performance',
                 [(fmt_pct(c['nm']), pc(c['nm'])), (fmt_pct(c['qtd_n']*100), pc(c['qtd_n'])),
                  (fmt_pct(c['ytd_n']*100), pc(c['ytd_n']))], y, bold=True)
    y -= 0.18

    # Since Inception
    y = draw_header(['SINCE INCEPTION','Cumulative Gross','Cumulative Net','Annualized Net'], y)
    y = draw_row(f"Since {fmt_mo(c['first_mo'])}",
                 [(fmt_pct(c['si_g']*100), pc(c['si_g'])),
                  (fmt_pct(c['si_n']*100), pc(c['si_n'])),
                  (fmt_pct(c['ann_n']*100), pc(c['ann_n']))], y)

    plt.subplots_adjust(left=0, right=1, top=1, bottom=0)
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0.05,
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════
# PPTX MODIFICATION: SLIDE 14 (FACT SHEET IMAGE)
# ═══════════════════════════════════════════════════════════════════
def gen_fact_image(client_name, strat_name, month, data, idx_ret):
    """Render fact sheet as a high-res PNG."""
    import matplotlib; matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    from matplotlib.patches import Rectangle

    si_d = [d for d in data if d['month_end_date'][:7] <= month]
    md = next((d for d in si_d if d['month_end_date'][:7]==month), None)
    if not md: return None
    year, n = month[:4], len(si_d)
    fm = si_d[0]['month_end_date'][:7] if si_d else month
    GOLD, DK, GRN, RED = '#CAA879','#1A1A2E','#059669','#DC2626'
    def pc(v): return GRN if v>=0 else RED

    rets = [(d.get('net_return_pct') or 0) for d in si_d]
    avg_r = sum(rets)/n if n else 0
    var = sum((r-avg_r)**2 for r in rets)/(n-1) if n>1 else 0
    ann_std = math.sqrt(var)*math.sqrt(12)
    pos_mo = sum(1 for r in rets if r>0)

    bench_id = 1
    b_rets = [(idx_ret.get(bench_id,{}).get(d['month_end_date'][:7]) or 0) for d in si_d]
    avg_b = sum(b_rets)/n if n else 0
    beta = corr = alpha = 0
    if n>2:
        cov = sum((rets[i]-avg_r)*(b_rets[i]-avg_b) for i in range(n))
        vb = sum((b-avg_b)**2 for b in b_rets)
        va = sum((r-avg_r)**2 for r in rets)
        beta = cov/vb if vb else 0
        corr = cov/math.sqrt(va*vb) if va and vb else 0
    alpha = (avg_r - beta*avg_b)*12

    def chain(lst): 
        c=1
        for d in lst: c *= (1+(d.get('net_return_pct') or 0))
        return c-1

    ytd_n = chain([d for d in si_d if d['month_end_date'][:4]==year])
    a1y = chain(si_d[-12:])
    a3y = None
    if len(si_d)>=12:
        l36 = si_d[-36:]
        c36 = chain(l36)
        a3y = math.pow(1+c36, 12/len(l36))-1
    si_n = chain(si_d)
    si_ann = (math.pow(1+si_n, 12/n)-1) if n>12 else si_n

    # VAMI
    vami = [100]
    for d in si_d: vami.append(vami[-1]*(1+(d.get('net_return_pct') or 0)))
    vami_lbl = [fm] + [d['month_end_date'][:7] for d in si_d]

    fig = plt.figure(figsize=(8.54, 5.2)); fig.patch.set_facecolor('white')

    # Header
    axi = fig.add_axes([0.02, 0.90, 0.96, 0.10]); axi.axis('off')
    axi.text(0, 0.7, client_name, fontsize=12, fontweight='bold', color=DK, va='top')
    axi.text(0, 0.15, strat_name, fontsize=8, fontweight='bold', color=GOLD, va='top')
    axi.text(0.45, 0.15, f'Inception: {fmt_mo(fm)}', fontsize=7, color='#999', va='top')

    # VAMI chart
    axv = fig.add_axes([0.03, 0.44, 0.44, 0.42])
    axv.plot(range(len(vami)), vami, color=GOLD, lw=1.5)
    axv.set_title('VAMI (Growth of $100)', fontsize=7, color=DK, loc='left', pad=3)
    step = max(1, len(vami)//8)
    axv.set_xticks(range(0, len(vami), step))
    axv.set_xticklabels([vami_lbl[i][-5:] for i in range(0, len(vami), step)], rotation=45, fontsize=4.5)
    axv.tick_params(axis='y', labelsize=5); axv.grid(axis='y', alpha=0.2, lw=0.4)
    axv.spines['top'].set_visible(False); axv.spines['right'].set_visible(False)

    # Risk metrics (right)
    axr = fig.add_axes([0.52, 0.56, 0.46, 0.32]); axr.axis('off')
    axr.add_patch(Rectangle((0, 0.87), 1, 0.13, fc=GOLD, ec='none'))
    axr.text(0.02, 0.93, 'RISK METRICS', fontsize=6.5, fontweight='bold', color='white', va='center')
    axr.text(0.58, 0.93, 'Gridiron', fontsize=6.5, fontweight='bold', color='white', va='center', ha='right')
    axr.text(0.88, 0.93, 'US Agg', fontsize=6.5, fontweight='bold', color='white', va='center', ha='right')
    rr = [('Std Deviation', fmt_pct(ann_std*100), '—'),
          ('Alpha', '', fmt_pct(alpha*100)),
          ('Beta', '', f'{beta:.2f}'),
          ('Correlation', '', f'{corr:.2f}'),
          ('% Positive Mo.', f'{pos_mo/n*100:.0f}%' if n else '—', '—')]
    for i,(l,g,b) in enumerate(rr):
        yp=0.74-i*0.15
        axr.text(0.02,yp,l,fontsize=6.5,color='#333',va='center')
        axr.text(0.58,yp,g,fontsize=6.5,color='#333',va='center',ha='right')
        axr.text(0.88,yp,b,fontsize=6.5,color='#666',va='center',ha='right')
        axr.plot([0,1],[yp-0.06,yp-0.06],color='#e0e0e0',lw=0.3)

    # Performance (right lower)
    axp = fig.add_axes([0.52, 0.14, 0.46, 0.38]); axp.axis('off')
    axp.add_patch(Rectangle((0, 0.88), 1, 0.12, fc=GOLD, ec='none'))
    axp.text(0.02, 0.94, 'PERFORMANCE', fontsize=6.5, fontweight='bold', color='white', va='center')
    axp.text(0.58, 0.94, 'Gridiron', fontsize=6.5, fontweight='bold', color='white', va='center', ha='right')
    axp.text(0.88, 0.94, 'US Agg', fontsize=6.5, fontweight='bold', color='white', va='center', ha='right')
    pp = [('MTD', (md.get('net_return_pct') or 0)*100), ('YTD', ytd_n*100),
          ('1Y Ann.', a1y*100), ('3Y Ann.', a3y*100 if a3y is not None else None),
          ('SI Cum.', si_n*100), ('SI Ann.', si_ann*100)]
    for i,(l,g) in enumerate(pp):
        yp=0.76-i*0.12
        axp.text(0.02,yp,l,fontsize=6.5,color='#333',va='center')
        if g is not None:
            axp.text(0.58,yp,fmt_pct(g),fontsize=6.5,color=pc(g),va='center',ha='right')
        else:
            axp.text(0.58,yp,'—',fontsize=6.5,color='#999',va='center',ha='right')
        axp.text(0.88,yp,'—',fontsize=6.5,color='#999',va='center',ha='right')
        axp.plot([0,1],[yp-0.05,yp-0.05],color='#e0e0e0',lw=0.3)

    # Monthly grid
    axg = fig.add_axes([0.02, 0.01, 0.96, 0.38]); axg.axis('off')
    years = sorted(set(d['month_end_date'][:4] for d in si_d))
    ms = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec','YTD']
    cw = 1.0/14
    axg.add_patch(Rectangle((0, 0.90), 1, 0.10, fc=GOLD, ec='none'))
    axg.text(cw*0.5, 0.95, 'Year', fontsize=4.5, fontweight='bold', color='white', va='center', ha='center')
    for mi, m in enumerate(ms):
        axg.text(cw*(mi+1.5), 0.95, m, fontsize=4.5, fontweight='bold', color='white', va='center', ha='center')

    rh = min(0.11, 0.88/max(len(years),1))
    for yi, yr in enumerate(years):
        yp = 0.84 - yi*rh
        axg.text(cw*0.5, yp, yr, fontsize=4.5, fontweight='bold', color=DK, va='center', ha='center')
        yc = 1
        for mi in range(12):
            ym = f"{yr}-{mi+1:02d}"
            rec = next((d for d in si_d if d['month_end_date'][:7]==ym), None)
            if rec and rec.get('net_return_pct') is not None:
                v = rec['net_return_pct']; yc *= (1+v)
                vd = v*100
                axg.text(cw*(mi+1.5), yp, f'{vd:.1f}%', fontsize=3.5, color=pc(vd), va='center', ha='center')
            else:
                axg.text(cw*(mi+1.5), yp, '-', fontsize=3.5, color='#ccc', va='center', ha='center')
        ytdv = (yc-1)*100
        axg.text(cw*13.5, yp, f'{ytdv:.1f}%', fontsize=3.5, fontweight='bold', color=pc(ytdv), va='center', ha='center')
        axg.plot([0,1],[yp-rh*0.4, yp-rh*0.4], color='#eee', lw=0.2)

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0.05,
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════
# IMAGE REPLACEMENT HELPER
# ═══════════════════════════════════════════════════════════════════
def replace_content_image(slide, img_buf):
    for s in slide.shapes:
        if s.shape_type == 13 and s.name == 'Content Image':
            l, t, w, h = s.left, s.top, s.width, s.height
            s._element.getparent().remove(s._element)
            slide.shapes.add_picture(img_buf, l, t, w, h)
            return True
    return False

def update_datebox(slide, month):
    for s in slide.shapes:
        if s.has_text_frame and s.name == 'DateBox':
            for p in s.text_frame.paragraphs:
                for r in p.runs: r.text = fmt_mo(month)
            return True
    return False


# ═══════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════
def main():
    ap = argparse.ArgumentParser(description='Gridiron Presentation Generator')
    ap.add_argument('--client', help='Client name (partial match)')
    ap.add_argument('--month', help='YYYY-MM')
    ap.add_argument('--template', default='test_slides12_14_v9_2.pptx')
    ap.add_argument('--output', help='Output path')
    ap.add_argument('--email', help='Supabase email')
    ap.add_argument('--password', help='Supabase password')
    ap.add_argument('--skip-auth', action='store_true', help='Use anon key (no login)')
    args = ap.parse_args()

    print("═" * 50)
    print("  GRIDIRON PRESENTATION GENERATOR")
    print("═" * 50)

    # Auth
    token = SB_KEY
    if args.skip_auth:
        print("\n  Using anon key (no login)")
    elif args.email and args.password:
        print("\nAuthenticating...")
        token = sb_login(args.email, args.password)
    else:
        e = input("\nEmail (or press Enter to skip auth): ").strip()
        if e:
            p = getpass.getpass("Password: ")
            if p: token = sb_login(e, p)
        else:
            print("  Using anon key (no login)")
    print("  ✓ Ready")

    # Load
    print("\nLoading data...")
    D = load_all(token)

    # Client
    if args.client:
        matches = [c for c in D['clients'] if args.client.lower() in c['client_name'].lower()]
        if not matches:
            print(f"No client matching '{args.client}'"); return
        client = matches[0]
    else:
        print("\nClients:")
        for i, c in enumerate(D['clients']):
            print(f"  {i+1:2d}. {c['client_name']}")
        client = D['clients'][int(input("\nSelect #: "))-1]
    print(f"\n→ Client: {client['client_name']}")

    mdata = load_monthly(client['client_id'], token)
    print(f"  {len(mdata)} months of data")
    if not mdata: print("  ERROR: No data!"); return

    # Month
    months = sorted(set(d['month_end_date'][:7] for d in mdata), reverse=True)
    if args.month:
        month = args.month
    else:
        print("\nMonths:")
        for i, m in enumerate(months[:12]):
            print(f"  {i+1:2d}. {fmt_mo(m)}")
        month = months[int(input("\nSelect #: "))-1]
    print(f"→ Month: {fmt_mo(month)}")

    strat = next((s for s in D['strategies'] if s['strategy_id']==client.get('strategy_id')), None)
    sname = strat['strategy_name'] if strat else 'Gridiron Tactical Fixed Income'

    # Sectors
    print("\nBuilding sector data...")
    sec_data = build_sectors(D['tickers'], D['cef_daily'])

    # Template
    tpath = args.template
    for p in [tpath, f'/mnt/user-data/uploads/{tpath}', os.path.expanduser(f'~/{tpath}')]:
        if os.path.exists(p): tpath = p; break
    print(f"\nTemplate: {tpath}")
    prs = Presentation(tpath)

    # ── Slide 1 ──
    print("\n[1/5] Cover...")
    update_cover(prs.slides[0], client['client_name'], month)
    print("    ✓ Client name & date")

    # ── Slides 7-11 ──
    print("[2/5] Sector discounts...")
    update_sector_slides(prs, sec_data, month)

    # ── Slide 12 ──
    print("[3/5] Asset allocation...")
    update_allocation(prs.slides[11], month)
    print("    ✓ Date header updated")

    # ── Slide 13 ──
    print("[4/5] Performance statement...")
    img = gen_perf_image(client['client_name'], sname, month, mdata)
    if img:
        update_datebox(prs.slides[12], month)
        replace_content_image(prs.slides[12], img)
        print("    ✓ Image generated & placed")

    # ── Slide 14 ──
    print("[5/5] Fact sheet...")
    img = gen_fact_image(client['client_name'], sname, month, mdata, D['index_returns'])
    if img:
        update_datebox(prs.slides[13], month)
        replace_content_image(prs.slides[13], img)
        print("    ✓ Image generated & placed")

    # Save
    safe = client['client_name'].replace("'","").replace("/","-").replace(" ","_")
    out = args.output or f"{safe}_{month}.pptx"
    prs.save(out)
    print(f"\n{'═'*50}")
    print(f"  SAVED: {out}")
    print(f"  Size:  {os.path.getsize(out)/1024:.0f} KB")
    print(f"{'═'*50}")
    print(f"\nOpen in PowerPoint → File → Save As → PDF")

if __name__ == '__main__':
    main()
